
"""
external_sources.py — URL/업로드 문서/이미지 텍스트 추출
지원: URL, TXT, CSV, XLSX/XLS, PDF, DOCX, PPTX, HWP/HWPX, PNG/JPG/WEBP 이미지 OCR optional
"""
from __future__ import annotations

import io
import re
import zipfile
from html import unescape
from pathlib import Path
from typing import Iterable

import requests
from bs4 import BeautifulSoup


def clean_text(text: str, limit: int = 50000) -> str:
    text = unescape(text or "")
    text = re.sub(r"\s+", " ", text).strip()
    return text[:limit]


def fetch_url_text(url: str, max_chars: int = 30000) -> dict:
    url = url.strip()
    if not url:
        return {"source": "URL", "title": "", "text": "", "error": "빈 URL"}
    try:
        headers = {"User-Agent": "Mozilla/5.0 (compatible; StockBriefingRAG/1.0)"}
        resp = requests.get(url, headers=headers, timeout=12)
        resp.raise_for_status()
        ctype = resp.headers.get("content-type", "")
        if "text" not in ctype and "html" not in ctype and "json" not in ctype:
            return {"source": url, "title": url, "text": clean_text(resp.text, max_chars)}
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "noscript", "svg", "header", "footer", "nav"]):
            tag.decompose()
        title = soup.title.get_text(" ", strip=True) if soup.title else url
        paragraphs = [p.get_text(" ", strip=True) for p in soup.find_all(["h1", "h2", "h3", "p", "li"])]
        text = clean_text("\n".join(paragraphs) or soup.get_text(" ", strip=True), max_chars)
        return {"source": url, "title": title, "text": text, "error": ""}
    except Exception as e:
        return {"source": url, "title": url, "text": "", "error": f"URL 읽기 실패: {type(e).__name__}: {e}"}


def _decode_bytes(data: bytes) -> str:
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            return data.decode(enc)
        except Exception:
            continue
    return data.decode("utf-8", errors="ignore")


def extract_txt(data: bytes) -> str:
    return _decode_bytes(data)


def extract_excel(data: bytes, suffix: str) -> str:
    import pandas as pd
    bio = io.BytesIO(data)
    try:
        xls = pd.ExcelFile(bio)
    except Exception:
        return "[엑셀 읽기 실패] openpyxl/xlrd 설치 또는 파일 형식을 확인하세요."
    parts = []
    for sheet in xls.sheet_names[:10]:
        try:
            df = xls.parse(sheet).head(80)
            parts.append(f"[Sheet: {sheet}]\n" + df.to_csv(index=False))
        except Exception as e:
            parts.append(f"[Sheet: {sheet}] 읽기 실패: {e}")
    return "\n\n".join(parts)


def extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    texts = []
    for i, page in enumerate(reader.pages[:30]):
        try:
            texts.append(f"[PDF page {i+1}]\n" + (page.extract_text() or ""))
        except Exception:
            continue
    return "\n\n".join(texts)


def extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for t in doc.tables:
        for row in t.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def extract_hwpx(data: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            names = [n for n in z.namelist() if n.endswith(".xml")]
            texts = []
            for name in names[:80]:
                raw = z.read(name).decode("utf-8", errors="ignore")
                soup = BeautifulSoup(raw, "xml")
                t = soup.get_text(" ", strip=True)
                if t:
                    texts.append(t)
            return clean_text("\n".join(texts), 50000)
    except Exception as e:
        return f"[HWPX 읽기 실패] {e}"


def extract_hwp(data: bytes) -> str:
    # HWP 구버전은 복잡하므로 먼저 미리보기 텍스트(PrvText) 스트림을 시도한다.
    try:
        import olefile
        ole = olefile.OleFileIO(io.BytesIO(data))
        for stream_name in ["PrvText", "\x05HwpSummaryInformation"]:
            if ole.exists(stream_name):
                raw = ole.openstream(stream_name).read()
                text = _decode_bytes(raw)
                if text.strip():
                    return clean_text(text, 50000)
        return "[HWP 읽기 제한] 이 HWP 파일에는 추출 가능한 PrvText가 없습니다. HWPX/DOCX/PDF로 변환해 업로드하면 정확도가 높습니다."
    except Exception as e:
        return f"[HWP 읽기 실패] {e}"


def extract_image_text(data: bytes, filename: str) -> str:
    try:
        from PIL import Image
        import pytesseract  # optional. requirements에는 넣지 않음. 설치되어 있으면 사용.
        img = Image.open(io.BytesIO(data))
        text = pytesseract.image_to_string(img, lang="kor+eng")
        return text.strip() or f"[이미지 OCR 결과 없음] {filename}"
    except Exception as e:
        return f"[이미지 파일 인식 제한] {filename}: OCR 엔진이 없거나 이미지 텍스트 추출에 실패했습니다. 필요 시 이미지의 주요 텍스트를 별도 입력란에 붙여주세요. ({type(e).__name__})"


def extract_uploaded_file(uploaded_file, max_chars: int = 50000) -> dict:
    name = uploaded_file.name
    suffix = Path(name).suffix.lower().lstrip(".")
    data = uploaded_file.getvalue()
    try:
        if suffix in {"txt", "md", "csv"}:
            text = extract_txt(data)
        elif suffix in {"xlsx", "xls"}:
            text = extract_excel(data, suffix)
        elif suffix == "pdf":
            text = extract_pdf(data)
        elif suffix in {"docx", "doc"}:
            text = extract_docx(data) if suffix == "docx" else "[DOC 읽기 제한] .doc 파일은 .docx 또는 PDF로 변환해 업로드하세요."
        elif suffix == "pptx":
            text = extract_pptx(data)
        elif suffix == "hwp":
            text = extract_hwp(data)
        elif suffix == "hwpx":
            text = extract_hwpx(data)
        elif suffix in {"png", "jpg", "jpeg", "webp", "bmp"}:
            text = extract_image_text(data, name)
        else:
            text = f"[지원하지 않는 파일 형식] {suffix}"
        return {"source": name, "title": name, "text": clean_text(text, max_chars), "error": ""}
    except Exception as e:
        return {"source": name, "title": name, "text": "", "error": f"파일 읽기 실패: {type(e).__name__}: {e}"}


def collect_external_sources(urls: Iterable[str], uploaded_files: Iterable) -> list[dict]:
    docs = []
    for url in urls:
        url = url.strip()
        if url:
            docs.append(fetch_url_text(url))
    for f in uploaded_files or []:
        docs.append(extract_uploaded_file(f))
    return docs
