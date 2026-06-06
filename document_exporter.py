
"""
document_exporter.py — 최종 리포트 Word 다운로드, 요약 PPT 다운로드 생성
"""
from __future__ import annotations

import io
import re
from datetime import date


def _strip_md(line: str) -> str:
    line = re.sub(r"[#*_`>]", "", line)
    return line.strip()


def make_docx_report(report_text: str, title: str = "주식 AI 브리핑 최종 리포트") -> bytes:
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    styles = doc.styles
    styles["Normal"].font.name = "Malgun Gothic"
    styles["Normal"].font.size = Pt(10.5)
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"생성일: {date.today().isoformat()}")
    for raw in report_text.splitlines():
        line = raw.rstrip()
        if not line:
            continue
        clean = _strip_md(line)
        if line.startswith("#") or line.startswith("【"):
            doc.add_heading(clean, level=1 if line.startswith("#") or line.startswith("【") else 2)
        elif line.startswith("-"):
            doc.add_paragraph(clean.lstrip("- "), style="List Bullet")
        else:
            doc.add_paragraph(clean)
    doc.add_paragraph("투자 참고용 자료이며 최종 투자 판단은 사용자 책임입니다.")
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _extract_summary_points(report_text: str, max_points: int = 18) -> list[str]:
    lines = []
    for raw in report_text.splitlines():
        clean = _strip_md(raw)
        if not clean or len(clean) < 8:
            continue
        if any(k in clean for k in ["요약", "전략", "예상", "리스크", "UP", "DOWN", "NEUTRAL", "매수", "관망", "손절", "익절"]):
            lines.append(clean[:130])
        if len(lines) >= max_points:
            break
    return lines or [_strip_md(x)[:130] for x in report_text.splitlines() if _strip_md(x)][:max_points]


def make_ppt_summary(report_text: str, stocks: list[dict] | None = None, title: str = "주식 AI 브리핑 요약") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"생성일: {date.today().isoformat()}\nAI 분석 기반 투자 브리핑 요약"

    def add_bullet_slide(slide_title: str, bullets: list[str]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets[:7]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)

    add_title_slide()
    if stocks:
        rows = []
        for s in stocks[:8]:
            metrics = s.get("metrics", {})
            rows.append(f"{s.get('name', s.get('ticker'))}: {s.get('close','N/A')} / {s.get('change_rate','N/A')}% / {metrics.get('heuristic_direction','N/A')}")
        add_bullet_slide("종목 현황 요약", rows)
    points = _extract_summary_points(report_text, 21)
    add_bullet_slide("핵심 투자전략", points[:7])
    add_bullet_slide("내일 대응전략", points[7:14])
    add_bullet_slide("리스크 체크리스트", points[14:21])
    add_bullet_slide("면책", ["본 자료는 AI 분석 기반 투자 참고용입니다.", "최종 투자 판단과 책임은 사용자에게 있습니다."])
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()
